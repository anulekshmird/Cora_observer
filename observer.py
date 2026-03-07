import time
import mss
import ollama
import threading
from PIL import Image
import io
import os
import config
import json
import re
import context_engine
import ocr_engine
from PyQt6.QtCore import QObject, pyqtSignal
import base64
import requests
import pyautogui
import hashlib
from datetime import datetime
import docx
from pptx import Presentation

class ObserverSignal(QObject):
    suggestion_ready = pyqtSignal(object) # json payload
    prepare_capture = pyqtSignal()
    finished_capture = pyqtSignal()
    error_resolved = pyqtSignal()

class Observer:
    def __init__(self):
        self.running = False
        self.paused = False
        self.stop_flag = False
        self.signals = ObserverSignal()
        self.model = config.OLLAMA_MODEL 
        self.context_engine = context_engine.ContextEngine()
        self.last_llm_call_time = 0
        
        # Proactive context storage (for grounded suggestion execution)
        self.last_ocr_text = ""
        self.last_proactive_screenshot = None  # bytes
        self.last_frame_hash = None # For Fix 1 (Deduplication)
        self.last_screen_hash = None # For proactive loop deduplication
        self.proactive_pause = False # For pausing proactive analysis
        self.last_reported_error_sig = None # For syntax error tracking
        
        # Session Management
        self.chats_dir = os.path.join(os.getcwd(), "chats")
        if not os.path.exists(self.chats_dir):
            os.makedirs(self.chats_dir)
            
        self.current_session_id = None
        self.chat_history = [] 
        self.create_new_session()

    def create_new_session(self):
        import uuid
        self.current_session_id = str(uuid.uuid4())[:8]
        self.chat_history = []
        print(f"Created new session: {self.current_session_id}")
        self.save_session()

    def switch_session(self, session_id):
        filepath = os.path.join(self.chats_dir, f"{session_id}.json")
        if os.path.exists(filepath):
            try:
                with open(filepath, 'r') as f:
                    data = json.load(f)
                    self.chat_history = data.get('history', [])
                self.current_session_id = session_id
                print(f"Switched to session: {session_id}")
                return True
            except Exception as e:
                print(f"Error loading session: {e}")
                return False
        return False

    def get_sessions(self):
        sessions = []
        if not os.path.exists(self.chats_dir): return []
        
        for f in os.listdir(self.chats_dir):
            if f.endswith(".json"):
                 sid = f.replace(".json", "")
                 # Load first message as title if poss?
                 title = f"Chat {sid}"
                 try:
                     with open(os.path.join(self.chats_dir, f), 'r') as file:
                         data = json.load(file)
                         
                         # Priority 1: Saved Title
                         if data.get('title'):
                             title = data['title']
                         else:
                             # Priority 2: First Message Inference (Fallback)
                             hist = data.get('history', [])
                             if hist:
                                 for msg in hist:
                                     if msg['role'] == 'user':
                                         txt = msg['content'].split("USER:")[-1].strip()[:30]
                                         title = txt if txt else title
                                         break
                 except: pass
                 sessions.append({'id': sid, 'title': title})
        return sessions

    def delete_session(self, session_id):
        filepath = os.path.join(self.chats_dir, f"{session_id}.json")
        try:
            if os.path.exists(filepath):
                os.remove(filepath)
                print(f"Deleted session: {session_id}")
                
                # If current session deleted, create new one
                if self.current_session_id == session_id:
                    self.create_new_session()
                return True
        except Exception as e:
            print(f"Error deleting session: {e}")
        return False

    def save_session(self):
        if not self.current_session_id: return
        filepath = os.path.join(self.chats_dir, f"{self.current_session_id}.json")
        try:
            # Load existing to preserve title
            # Strip image bytes from history (not JSON serializable)
            clean_history = []
            for msg in self.chat_history:
                clean_msg = {k: v for k, v in msg.items() if k != 'images'}
                clean_history.append(clean_msg)
            
            data = {'id': self.current_session_id, 'history': clean_history}
            if os.path.exists(filepath):
                 with open(filepath, 'r') as f:
                     existing = json.load(f)
                     if 'title' in existing:
                         data['title'] = existing['title']
            
            with open(filepath, 'w') as f:
                json.dump(data, f, indent=2)
        except Exception as e:
            print(f"Error saving session: {e}")

    def stop_chat(self):
        self.stop_flag = True
        print("Stopping generation...")

    def clear_history(self):
        # Instead of clearing, we create a new session
        self.create_new_session()

    # ... (capture_screen, _image_to_bytes, pause, resume, analyze, read_file_content unused changes omitted)



    def capture_screen(self):
        try:
            # 0. Prevent Self-Analysis (Recursion Guard)
            win_title = self.context_engine.get_active_window_title().lower()
            if any(x in win_title for x in ["cora", "assistant", "suggestion", "overlay"]):
                return None

            # 1. Hide UI (Prevent recursion)
            self.signals.prepare_capture.emit()
            time.sleep(0.3) # Give UI time to vanish
            
            # FIX 7: Try Partial Screen Capture for Developer Mode
            mode = self.context_engine.get_context_snapshot().get('mode_primary', 'general')
            
            import mss
            with mss.mss() as sct:
                monitor = sct.monitors[1] # Primary
                region = monitor # Full screen default
                
                # Active Window Capture for Developer Mode
                if mode == "developer":
                    try:
                        import pygetwindow as gw
                        win = gw.getActiveWindow()
                        if win and win.width > 0 and win.height > 0:
                            # Constrain to primary monitor boundaries
                            left = max(win.left, monitor["left"])
                            top = max(win.top, monitor["top"])
                            right = min(win.left + win.width, monitor["left"] + monitor["width"])
                            bottom = min(win.top + win.height, monitor["top"] + monitor["height"])
                            
                            region = {
                                "top": int(top),
                                "left": int(left),
                                "width": int(right - left),
                                "height": int(bottom - top)
                            }
                            print(f"Observer: Capturing Window ({region['width']}x{region['height']})")
                    except Exception as e:
                        print(f"Window capture fallback to full: {e}")

                sct_img = sct.grab(region)
                img = Image.frombytes("RGB", sct_img.size, sct_img.bgra, "raw", "BGRX")
                # Downscale for performance, but KEEP READABLE
                img.thumbnail((3000, 3000))
                
            # 2. Restore UI
            self.signals.finished_capture.emit()
            return img
            
        except Exception as e:
            print(f"Screen Capture Error: {e}")
            self.signals.finished_capture.emit() # Always restore
            return None

    def _image_to_bytes(self, image):
        if not image: return None
        with io.BytesIO() as output:
            image.save(output, format='PNG') # PNG is lossless, better for text
            return output.getvalue()

    def pause(self):
        self.paused = True
        print("Observer Paused for Chat.")

    def resume(self):
        self.paused = False
        print("Observer Resumed.")

    def extract_text_from_screen(self, image):
        """
        Extracts visible text from the screen image using OCR.
        """
        from ocr_engine import extract_text
        return extract_text(image)

    def analyze(self, image, context_text=""):
        if self.paused or not image: return None
        
        ocr_text = ""
        # FIX 8: Self-analysis guard extension
        try:
            win_title = self.context_engine.get_active_window_title().lower()
            # IGNORE SYSTEM WINDOWS (Task Switching, Start Menu, etc.)
            system_titles = ["task switching", "task view", "start", "search", "notification center", "action center", "new notification", "cortana", "volume control", "system tray", "windows shell", "microsoft shell"]
            if any(kw in win_title for kw in ["cora", "assistant", "suggestion", "overlay"]) or \
               any(t == win_title or t in win_title for t in system_titles) or \
               not win_title.strip() or win_title == "window":
                return None
        except:
            pass
        
        # Convert to bytes if PIL Image
        image_data = None
        if image:
            image_data = self._image_to_bytes(image)
        
        if image_data is None:
            return None
        
        # FIX 1: Screen Hash Deduplication
        import hashlib
        current_hash = hashlib.md5(image_data).hexdigest()
        if current_hash == self.last_frame_hash:
             print("Observer: Screen hash match. Skipping redundant analysis.")
             return None
        self.last_frame_hash = current_hash

        # -----------------------------------------------------------------
        # HYBRID PERCEPTION: OCR + VISION
        # -----------------------------------------------------------------
        ocr_text = ""
        try:
             # FIX 2 & 3: Conditional OCR Logic
             snapshot = self.context_engine.get_context_snapshot()
             mode_primary = snapshot.get('mode_primary', 'general')
             win_title = snapshot.get('window_title', '').lower()
             
             # OCR Required?
             high_text_apps = ["word", "pdf", "docs", "notepad", "editor", ".pdf", "powerpoint", "slides", "keynote", "prezi"]
             need_ocr = mode_primary in ["developer", "writing", "reading"] or \
                        any(a in win_title for a in high_text_apps)
             
             if need_ocr:
                 # Re-convert bytes back to PIL for OCR
                 ocr_img = Image.open(io.BytesIO(image_data))
                 
                 # FIX 4: OCR Image Downscale (Max 1500px width for speed)
                 if ocr_img.width > 1500:
                      h = int(ocr_img.height * (1500 / ocr_img.width))
                      ocr_img = ocr_img.resize((1500, h), Image.Resampling.LANCZOS)
                 
                 ocr_text = ocr_engine.extract_text(ocr_img)
                 if len(ocr_text) < 20: 
                     ocr_text = "" # Ignore noise
                 else:
                     ocr_text = ocr_text[:2000] # Truncate for prompt
             else:
                 print(f"Observer: Skipping OCR for {mode_primary} mode (Light Analysis).")
                 
        except Exception as e:
             print(f"OCR Pipeline Error: {e}")
        
        # Store for suggestion execution pipeline
        self.last_ocr_text = ocr_text
        self.last_proactive_screenshot = image_data

        # REFINED PROACTIVE PROMPT: Grounded, Concise, and Adaptive
        full_prompt = f"""
        You are a subtle screen assistant. 
        PRIMARY CONTEXT (OCR): {ocr_text}
        ATTACHED SCREEN: (image)
        WINDOW HINT: {context_text}

        TASK: Describe visible content and suggest actions.
        
        STRICT RULES:
        1. GROUNDING: Describe ONLY what is visible. Do NOT infer app name unless visible in screenshot.
        2. PRIORITY: OCR Text > Visual Cues > Window Hint.
        3. HALLUCINATION: If OCR clearly indicates a document/slide, ignore window hints. Do not mix app contexts.
        4. REASON (Short): Max 12 words. High impact.
        5. REASON_LONG: 1-2 sentences of detail/explanation if useful.
        6. PPT TUNING: If bullet/slide structure, suggest: Summarize slide, Explain point, Speaker notes, Key points.
        7. OUTPUT: No disclaimers or coaching.

        OUTPUT JSON:
        {{
         "reason": "Visible observation (≤12 words)",
         "reason_long": "Optional detailed explanation",
         "confidence": 0.0-1.0, 
         "suggestions": [{{"label": "Summarize Slide", "hint": "Detail"}}]
        }}
        """
        
        try:
            # Rate Limiting
            now = time.time()
            if now - self.last_llm_call_time < 1.5:
                return None
            self.last_llm_call_time = now

            # Pick specialized prompt
            system_prompt = config.SYSTEM_PROMPT
            if mode_primary == 'developer':
                system_prompt = config.DEV_SYSTEM_PROMPT
            elif mode_primary == 'writing':
                system_prompt = config.PRODUCTIVITY_SYSTEM_PROMPT
            elif mode_primary == 'document':
                system_prompt = config.DOCUMENT_SYSTEM_PROMPT
            elif mode_primary == 'reading':
                system_prompt = config.READING_SYSTEM_PROMPT

            response = ollama.chat(model=self.model, messages=[
                {'role': 'system', 'content': system_prompt},
                {'role': 'user', 'content': full_prompt, 'images': [image_data]}
            ])
            text = response['message']['content'].strip()
            print(f"DEBUG: RAW OBSERVER OUT: {text[:100]}...") # Limit log

            # Clean JSON
            if "```json" in text:
                text = text.split("```json")[1].split("```")[0].strip()
            elif "```" in text:
                text = text.split("```")[1].split("```")[0].strip()
            
            # Loose JSON fix
            if not text.endswith("}"): 
                 idx = text.rfind("}")
                 if idx != -1: text = text[:idx+1]

            payload = json.loads(text)
            payload["screen_context"] = ocr_text
            return payload
        except Exception as e:
            # print(f"Observer Analyze Error: {e}")
            return None
            # print(f"Ollama Analyze Error: {e}") 
            return None

    def update_session_title(self, session_id, user_text):
        if not user_text: return
        try:
            # Generate a short 3-5 word title
            prompt = f"Summarize this user query into a short 3-5 word title: '{user_text}'. Return ONLY the title, no quotes."
            response = ollama.chat(model=self.model, messages=[
                {'role': 'user', 'content': prompt}
            ])
            title = response['message']['content'].strip().replace('"', '')
            
            # Save the new title
            filepath = os.path.join(self.chats_dir, f"{session_id}.json")
            if os.path.exists(filepath):
                with open(filepath, 'r+') as f:
                    data = json.load(f)
                    data['title'] = title
                    f.seek(0)
                    json.dump(data, f, indent=2)
                    f.truncate()
            print(f"Session {session_id} renamed to: {title}")
            return title
        except Exception as e:
            print(f"Title Generation Error: {e}")
            return None

    def read_pdf(self, path):
        try:
            import pypdf
            reader = pypdf.PdfReader(path)
            text = ""
            for page in reader.pages: # Remove limit to allow full document analysis
                 extract = page.extract_text()
                 if extract:
                     text += extract + "\n"
            
            if len(text.strip()) < 50:
                return "[WARNING: PDF appears to be scanned images or empty. Please ensure it contains selectable text.]"
                
            print(f"PDF Parsing Success: {len(text)} chars extracted.")
            return text
        except Exception as e:
            return f"[Error reading PDF: {e}]"

    def read_docx(self, path):
        try:
            doc = docx.Document(path)
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            return "\n".join(text)
        except Exception as e:
            return f"[Error reading DOCX: {e}]"

    def read_pptx(self, path):
        try:
            prs = Presentation(path)
            text = []
            for i, slide in enumerate(prs.slides):
                text.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            return f"[Error reading PPTX: {e}]"

    def read_file_content(self, path):
        try:
            if not path: return None
            ext = os.path.splitext(path)[1].lower()
            
            # 1. Specialized Parsers
            if ext == '.pdf':
                return self.read_pdf(path)
            elif ext == '.docx':
                return self.read_docx(path)
            elif ext == '.pptx':
                return self.read_pptx(path)

            # 2. Text/Code (Fallback)
            valid_exts = ['.txt', '.py', '.md', '.json', '.html', '.css', '.js', '.csv', '.bat', '.sh', '.xml', '.yaml', '.yml', '.ini', '.log']
            if ext not in valid_exts:
                return f"[File type '{ext}' not currently supported for deep analysis, but path is: {path}]"
            
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(100000) # Increased char limit for deep context
                return content
        except Exception as e:
            return f"[Error reading file: {e}]"

    def hash_screen(self, image):
        """Generates MD5 hash of screen image bytes to detect changes."""
        if image is None: return None
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        return hashlib.md5(img_byte_arr.getvalue()).hexdigest()

    def stream_chat_with_screen(self, user_query, attachment=None, proactive_context=None):
        self.stop_flag = False
        mode_primary = "general" # Default initialization
        try:
            image_bytes = None
            current_images = []
            prompt_context = ""
            
            # --- IMPROVEMENT 2: ATTACHMENT PRIORITY ---
            if attachment:
                print(f"Attachment detected: {attachment}. Skipping screen/OS context.")
                # Read attachment content
                content = self.read_file_content(attachment)
                
                # Format prompt context based on file type
                fname = os.path.basename(attachment)
                prompt_context = f"\n\n[PRIORITY ATTACHMENT: {fname}]\n\n{content}\n\n[END ATTACHMENT]\n"
                
                # Ensure attachment priority by skipping vision entirely
                mode_primary = "general"
                system_prompt = config.CHAT_SYSTEM_PROMPT
                # Skip vision logic and proceed to user content prep
            
            # --- IMPROVEMENT 1: REACTIVE VS PROACTIVE SEPARATION ---
            elif proactive_context:
                # Suggestion execution logic (proactive context provided)
                print("Using proactive context.")
                pc_mode = proactive_context.get('mode_primary', 'general')
                pc_window = proactive_context.get('window_title', 'Unknown')
                prompt_context = f"\n\n[COMMAND MODE: Suggestion Execution]\nACTIVE APP: {pc_window}\nMODE: {pc_mode}\n"
                
                if proactive_context.get('screenshot'):
                    current_images.append(proactive_context['screenshot'])
                
                mode_primary = pc_mode
                system_prompt = config.CHAT_SYSTEM_PROMPT # Or use mode-specific if needed
            
            else:
                # Normal chat - Check for vision keywords
                vision_keywords = ["look", "see", "screen", "visual", "watch", "what is this", "screenshot", "observe", "check", "debug", "fix"]
                is_vision_request = any(k in user_query.lower() for k in vision_keywords)
                
                if is_vision_request:
                    print("Vision keyword detected. Capturing screen.")
                    img = self.capture_screen()
                    if img:
                        image_bytes = self._image_to_bytes(img)
                        if image_bytes: 
                            current_images.append(image_bytes)
                            # --- IMPROVEMENT 3: OCR INJECTION ---
                            # Re-convert bytes back to PIL for OCR
                            ocr_img = Image.open(io.BytesIO(image_bytes))
                            ocr_text = ocr_engine.extract_text(ocr_img)
                            if ocr_text:
                                prompt_context = f"\n\nSCREEN_TEXT:\n{ocr_text[:2000]}\n"
                else:
                    print("Reactive Mode: Normal Conversation (No screen capture).")
                
                # Get OS snapshot for mode/context but NO screen capture
                os_context = self.context_engine.get_context_snapshot()
                window_title = os_context.get('window_title', 'Unknown')
                mode_primary = os_context.get('mode_primary', 'general')
                
                if is_vision_request:
                    prompt_context += f"\nWindow: {window_title}\n"
                
            # Select system prompt
            if mode_primary == 'developer':
                system_prompt = config.DEV_SYSTEM_PROMPT
            elif mode_primary == 'writing': 
                system_prompt = config.PRODUCTIVITY_SYSTEM_PROMPT
            elif mode_primary == 'document':
                system_prompt = config.DOCUMENT_SYSTEM_PROMPT
            elif mode_primary == 'reading':
                system_prompt = config.READING_SYSTEM_PROMPT
            elif mode_primary == 'video':
                system_prompt = config.VIDEO_SYSTEM_PROMPT
            else:
                system_prompt = config.CHAT_SYSTEM_PROMPT
            
            print(f"Streaming ({self.model})...")
            
            # 7. Construct History-Aware Message
            user_content = f"{prompt_context}\nUSER: {user_query}"
            
            new_message = {'role': 'user', 'content': user_content}
            
            # Ensure proper image handling for Ollama
            if current_images: 
                # Ollama expects list of base64 strings OR bytes.
                # Since _image_to_bytes returns bytes, and we read file as bytes, we are consistent.
                new_message['images'] = current_images
                
            self.chat_history.append(new_message)
            
            # Generate Title if First Message
            if len(self.chat_history) == 1:
                t = threading.Thread(target=self.update_session_title, args=(self.current_session_id, user_query), daemon=True)
                t.start()
            
            # 8. Send to LLM
            messages_payload = [{'role': 'system', 'content': system_prompt}] + self.chat_history
            stream = ollama.chat(model=self.model, messages=messages_payload, stream=True)

            full_response = ""
            for chunk in stream:
                if self.stop_flag: break
                token = chunk['message']['content']
                full_response += token
                yield token
            
            self.chat_history.append({'role': 'assistant', 'content': full_response})
            self.save_session()

        except Exception as e:
            print(f"Stream Error: {e}")
            yield f"[Error: {e}]"

    def _check_syntax_errors(self, ctx=None):
        # If ctx is not provided, get it from context_engine
        if ctx is None:
            ctx = self.context_engine.get_context_snapshot()

        if ctx.get('error'):
            sig = ctx['error_signature']
            if sig != self.last_reported_error_sig:
                # NEW ERROR DETECTED!
                print(f"🚨 New Syntax Error: {ctx['error']['message']} in {os.path.basename(ctx['error']['file'])}")

                # Generate Fix Suggestions via LLM (Silent)
                error_prompt = f"""
                SYNTAX ERROR DETECTED:
                File: {ctx['error']['file']}
                Line: {ctx['error']['line']}
                Error: {ctx['error']['message']}
                Code:
                {ctx['error']['context']}

                Provide a brief fix explanation and the corrected code block.
                Format as JSON: {{ "reason": "Explanation", "code": "Corrected Code", "confidence": 1.0 }}
                """

                # Call LLM
                response = ollama.chat(model=self.model, messages=[
                     {'role': 'system', 'content': config.DEV_SYSTEM_PROMPT},
                     {'role': 'user', 'content': error_prompt}
                ])

                # Parse
                text = response['message']['content'].strip()
                # Clean JSON
                if "```json" in text:
                    text = text.split("```json")[1].split("```")[0].strip()
                elif "```" in text:
                    text = text.split("```")[1].split("```")[0].strip()

                try:
                    payload = json.loads(text)
                    payload['type'] = 'syntax_error' # Mark for UI
                    payload.setdefault('screen_context', '')
                    payload.setdefault('error_context', '')
                    payload.setdefault('suggestions', [])
                    self.signals.suggestion_ready.emit(payload)
                    self.last_reported_error_sig = sig # Mark handled
                except:
                    pass

    def loop(self):
        print("Observer started (Silent Mode)...")
        self.running = True
        self.last_reported_error_sig = None
        self.last_screen_hash = None # Initialize last screen hash
        self.proactive_pause = False # Initialize proactive_pause
        
        while self.running:
            if self.paused or self.proactive_pause:
                time.sleep(1)
                continue

            try:
                # ----------------------------------------------
                # 1. Proactive Monitoring (Throttled by Hashing)
                # ----------------------------------------------
                
                # Capture Screen
                screenshot = self.capture_screen()
                if not screenshot:
                    time.sleep(1)
                    continue
                
                # Check Hash
                current_hash = self.hash_screen(screenshot)
                if current_hash == self.last_screen_hash:
                    # Screen unchanged, only check OS context for syntax errors
                    self._check_syntax_errors()
                    time.sleep(config.CHECK_INTERVAL) # Use config.CHECK_INTERVAL for consistency
                    continue
                
                self.last_screen_hash = current_hash
                
                # 2. OCR and Snapshot
                ocr_text = ocr_engine.extract_text(screenshot)
                ctx = self.context_engine.get_context_snapshot()
                
                # 3. Check for Syntax Errors (Always check when screen/context changes)
                self._check_syntax_errors(ctx)
                
                # 4. Visual Analysis (Only if not already handled by syntax or if in general/terminal)
                check_visual = True
                if ctx.get('mode_primary') == 'developer' and ctx.get('error'):
                    check_visual = False # Syntax error already takes precedence
                
                if check_visual:
                    # Add OCR text to context for visual analysis
                    context_text = ctx.get('window_title', '')
                    if ocr_text:
                        context_text += f"\nSCREEN_TEXT:\n{ocr_text[:2000]}" # Limit OCR text to avoid prompt overflow

                    payload = self.analyze(screenshot, context_text=context_text)
                    
                    if payload:
                        reason = payload.get('reason', '')
                        confidence = payload.get('confidence', 0.0)

                        # FILTER 1: Self-Reflection Prevention
                        if "Cora" in reason or "AI" in reason or "Ui" in reason:
                            pass # Skip
                        
                        # FILTER 2: Low Confidence Prevention
                        elif confidence < config.PROACTIVE_THRESHOLD:
                            pass # Skip low confidence
                        else:
                            self.signals.suggestion_ready.emit(payload)
                
                # self.loop_count += 1
            except Exception as e:
                print(f"Observer Loop Error: {e}")
            
            # Wait for next cycle
            time.sleep(config.CHECK_INTERVAL)
            time.sleep(config.CHECK_INTERVAL)

    def stop(self):
        self.running = False
